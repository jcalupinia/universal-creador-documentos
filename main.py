# ── Standard library ────────────────────────────────────────────────────────────
import base64
import io
import os

if os.name == "nt":
    for p in (r"C:\msys64\mingw64\bin",):
        if os.path.isdir(p):
            try:
                os.add_dll_directory(p)
            except Exception:
                pass

import random
import re
import tempfile
import urllib.request
import uuid
from datetime import date
from typing import Any, Dict, List, Optional, Union
from base64 import b64decode
from typing import List, Dict, Optional, Union, Any
from pydantic import BaseModel

# ── Third-party ────────────────────────────────────────────────────────────────
import matplotlib.pyplot as plt
import pandas as pd
from fastapi import FastAPI, File, HTTPException, UploadFile, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from jinja2 import Template
from pydantic import BaseModel
import json, ast
from typing import Any

# openpyxl (Excel)
from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.formatting.rule import ColorScaleRule
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter
from openpyxl.workbook.defined_name import DefinedName
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.worksheet.table import Table, TableStyleInfo
from openpyxl.drawing.image import Image as XLImage

# python-pptx (PowerPoint)
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

# python-docx (Word)
from docx import Document
from docx.enum.section import WD_ORIENT, WD_SECTION_START
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor

# PDF
from fpdf import FPDF
try:
    from weasyprint import HTML, CSS 
except Exception:
    HTML = None

# SVG/PNG
try:
    import cairosvg  
except Exception:
    cairosvg = None
    
def clean_text(text):
    if isinstance(text, str):
        return re.sub(r"[^\w\s\-.,()#]", "", text)
    return text

def _load_logo_bytes(logo_url: Optional[str] = None, logo_b64: Optional[str] = None) -> Optional[bytes]:
    """Devuelve los bytes del logo priorizando base64/url y con fallback corporativo."""
    if logo_b64:
        try:
            if logo_b64.startswith("data:"):
                logo_b64 = logo_b64.split(",", 1)[1]
            return b64decode(logo_b64)
        except Exception:
            pass
    url = logo_url or ""
    if url.startswith("data:"):
        try:
            return b64decode(url.split(",", 1)[1])
        except Exception:
            return b64decode(DEFAULT_LOGO_B64)
    if url.startswith("http"):
        try:
            with urllib.request.urlopen(url, timeout=10) as resp:
                return resp.read()
        except Exception:
            return b64decode(DEFAULT_LOGO_B64)
    if url:
        try:
            with open(url, "rb") as fh:
                return fh.read()
        except Exception:
            return b64decode(DEFAULT_LOGO_B64)
    return b64decode(DEFAULT_LOGO_B64)

def _logo_to_data_uri(logo_url: Optional[str] = None, logo_b64: Optional[str] = None) -> str:
    """Retorna el logo en formato data URI, con fallback corporativo."""
    try:
        data = _load_logo_bytes(logo_url, logo_b64)
        if not data:
            return DEFAULT_LOGO_DATA_URI
        return f"data:image/png;base64,{base64.b64encode(data).decode('ascii')}"
    except Exception:
        return DEFAULT_LOGO_DATA_URI

def _add_docx_image(run, image_bytes: bytes, width_in: float = 1.6):
    """Adjunta una imagen a un run de docx usando un archivo temporal."""
    if not image_bytes:
        return
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    try:
        tmp.write(image_bytes)
        tmp.flush()
        tmp.close()
        run.add_picture(tmp.name, width=DocxInches(width_in))
    finally:
        try:
            os.unlink(tmp.name)
        except Exception:
            pass

def _render_cover(doc, placeholders: dict, brand: dict):
    """Crea una portada simple centrada con título/subtítulo/autor/fecha."""
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt as DocxPt

    titulo = placeholders.get("titulo") or "Documento"
    subtitulo = placeholders.get("subtitulo") or ""
    autor = placeholders.get("autor") or ""
    fecha = placeholders.get("fecha") or ""

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(titulo)
    r.bold = True
    r.font.size = DocxPt(26)

    if subtitulo:
        p2 = doc.add_paragraph()
        p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r2 = p2.add_run(subtitulo)
        r2.font.size = DocxPt(14)

    meta = []
    if autor:
        meta.append(autor)
    if fecha:
        meta.append(fecha)
    if meta:
        p3 = doc.add_paragraph()
        p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r3 = p3.add_run(" – ".join(meta))
        r3.italic = True

    # salto de página tras portada
    doc.add_page_break()


def sanitize(data):
    if isinstance(data, dict):
        return {k: sanitize(v) for k, v in data.items()}
    elif isinstance(data, list):
        return [sanitize(i) for i in data]
    elif isinstance(data, str):
        return clean_text(data)
    return data

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

RESULT_DIR = "resultados"
os.makedirs(RESULT_DIR, exist_ok=True)
PUBLIC_BASE_URL = (os.getenv("PUBLIC_BASE_URL") or "").rstrip("/")
PDF_BASE_URL = (os.getenv("PDF_BASE_URL") or "https://universal-creador-documentos.onrender.com").rstrip("/")

def _result_url(filename: str, request: Optional[Request] = None) -> str:
    if PUBLIC_BASE_URL:
        return f"{PUBLIC_BASE_URL}/resultados/{filename}"
    if request is not None:
        try:
            return str(request.url_for("get_file", filename=filename))
        except Exception:
            pass
    return f"/resultados/{filename}"

def _pdf_url(filename: str) -> str:
    return f"{PDF_BASE_URL}/resultados/{filename}"

DEFAULT_COMPANY_NAME = "Audit Consulting Group"
DEFAULT_LOGO_URL = "https://i0.wp.com/auditconsulting.ec/wp-content/uploads/2023/02/Logo-color-Audit.png?fit=768%2C768&ssl=1"
DEFAULT_LOGO_B64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAwAAAAMACAMAAACkX/C8AAADAFBMVEVHcEwCAlIDA1GPuFQAAFEDA1GR"
    "qzwAAFEDA1EDA1EAAFKcukwAAFgDA1EEBE+myjQDA1EGBk4DA1EAAFIDA1EDA1EDA1ADA1EAAE8DA1AD"
    "A1EAAFADA1IFBVIAAFEAAFKYr0gDA1GYskgCAlEDA1EDA1ECAlAGBk4AAFADA1MDA1EDA1EDA1ADA1IC"
    "AlCUsEACAlEDA1EDA1EAAFACAlEDA1CZqUIDA1EDA1ADA1ADA1ECAlEAAFECAlGSrkICAlEDA1ECAlEC"
    "AkwDA1ECAlEEBFADA1GQpz8CAlEDA1EAAFIDA1GRqECZtEcDA1EDA1ADA1GZtUgDA1GKnDYEBFGOpDyO"
    "ojqTq0ACAlGYs0cCAlEDA1CVrUKNoDmTqUACAlEDA1ECAlGWr0WWr0SRpj6Rpz6PojqTq0GNoTuZtUmV"
    "sEWMnjgDA1Gbt0qUp0CbuEmVr0MDA1EDA1EDA1GbuEqZtEcDA1GNoTqSqD+OpTyXskWMnTeQpTwDA1AD"
    "A1GLnTeZtUgCAlGYs0eUq0GMnjiNnzmZtkqTpz+Qpj2Rpz4CAlGPozuPozuUq0ECAlEDA1GWr0SKmjQB"
    "AVGUrEGPpDyMnjiNoDmSqT8CAlGMnjiPozwDA1GMoDmSqUCZtkiUrEKVrkOMnjiUrEGTqT+Rpj4DA1GX"
    "sUUDA0+Uq0CTqkCTq0GMnziTqkACAlGMnjiNnzmUrEKMnjgCAlGbuEqMnjiPozuMnjgDA1ECAlECAlGP"
    "ozuatkmXsUWTq0GWr0SSqD8DA1GZtEeZtEiQpTyWr0SRpz6OoTmUrUKMnjibuEuRqD+bt0oAAFWatkkC"
    "Akybt0qbuUuYskabuUubuEuTqkGVrUKJmDMDA1GatkmYs0eat0qZtUibuUubuEubuEqYtEeWsEWLnTeR"
    "pj2MnjicukycuUuUrEKNoDmZtUmSqUCZtEiRpz6PozuQpT2VrkOXsUWOoTqOojuVrUOWr0SSqD+TqkGZ"
    "tkmTq0GNoDqXskaPpDyQpDyKmjSLmzaTqkCbt0otEQ28AAAA13RSTlMAZ1gSCOkJAfr+ENoEqyoB5iH4"
    "DuP9suwXOvYKfSgFH0Pzx3mc1XoaFEPwVY2tLAbdwJALlGMDpk1Zo84l1BefvHECNeBB9Ax1UBu3Dn8w"
    "7j3shkpGSj1KgfGZXJ/xGyPK0R9wn/Kk8Tqcjd7HoxNK8rBfSvFYxPfpIkrzdajY+ZTbn1PonncqMPpu"
    "rfI20opk8iK9XbiK4H/MJpFIl4jN6WT6KMXT4R1GiazTbGuEk9VUzMV67G9+aoPfrT3cprba9iWA/NT8"
    "esP9kbce2h7P9fr54oPDKo/CyWcAACAASURBVHja7N3LbxN3AsBxZ9dZO9mQ5uU8G6OIUEJ4LDSwlJIt"
    "iZOgPKqEVjmgIG1gb0hpJSjilqhc9sABpJ6Qeo44ceml5VbtueqiXoKqtqdW2v9iF8j8Uh52dgjYkzif"
    "z6mNPcZjz9fz+nmcSgEAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAABVoHHqDzvLocXXnueB1ow3nnWdJ8Zrd5S2a689z/smW+OmMty7rmeTO6Wj+/QOvHTL"
    "Yu9Lvu8ceqUnm2mOJpxueP6Gwt7eeIbkvonC8Ex2bSfJ1rzuLB+41LYYc5kYvdm0rqtQegntje4zPvri"
    "TcNLTS9rG7/dfuL+WNwM0sej6fo7n7+htb0pnomCxXzTAkayuyuAW7V1+9Lx7vp2X/SvbrIQZY5G92l5"
    "+6Xty5tFZ6Eu23JucHz5YmesAP4WTXVk+oUAlmK+Yp8JYPN1bO9SbjcFUH88tzbemWQA6xXksk39Y3kB"
    "bAOLl3K7KICatsePcqsh6QCe6m6f6MkIIPlDQWdyuyaA+vYn83qzdVsEsLbWt9CYF0DiZldzuyWA08ee"
    "Psy1bRLA2trtQ/UZASRtbH/d7gig9dT6Lv+Z9HYJYG1w33RBAEk7/+3uCGBqf/Q4UxUIoHe1u4jsi2vb"
    "3OHmwhYCmOmI56AA4ji4shsCaF3ojh5nJF3+AA58fLaI5Zml1bZzz83UwneZVw4gfW1fPIsCiHU+4L2V"
    "XRDA/JnwOC3z5Q+g1Es9MFbTf6npmRVB3eGBVw6ANyx/aKXqAxg63Lex3XE1nVQAT1/to5Mjtb/vd2VH"
    "8wJIWvr0SrUHcP7G7w/UcT7JAFKphsYLl1s2ns3chAASN3RypboD6Lk798wDvZ9PNIDHOyRTMxsrpLXj"
    "914zgPR081N7ByzJW95FrFmp6gDGRp59pLaxhANIFbqWBzfm60Lh9QJYfND/1MI1C/KW1d8/VsUBpCdr"
    "nxuHcKKQcACpzNhy+Bfqjs2+XgBd0QHebL/leOsGRo9VbwBHrz53tq9uZTbpAFKZrpmwH9BSvEcBVFbn"
    "5Hi1BpCvaXphHM6dQtIBpApT4Rx83c1GAWwDB95pq9IAhk9FR96b5tb/I/dVY+IBpIbudoQevxbAdjB9"
    "t6kqAyhcDGWf6I/Ow3Z8nXwAqS/CF5JyS60CKKu//iXONwEz9z6urcYAmvujFUDHwdnx9c2O3I3h5API"
    "fBNe7rYJAZTVh//8IFYB393pqL4AClNh9/796YaF6EO3djT5AFIHwjfyBh8IoKz++O8rn8b6TNr75WDV"
    "BTD9brTQn5vKp8aiwrMz08kHkLkwGJ5NQQBlDeDPDz/5MNan5d6FvmoLYD6c42t/vNmTPxJ96LbVJB9A"
    "qivsnNy+J4DyBvDrw0/+FOtDqflsX3UFMPAgWgFkrz0ZBDcRHX3PHhlIPoB02AY6Ni+AMgfw6w/xCijs"
    "Xe6uqgC6LocxN7NP9oPyx6P/vXkx+QBSJ6LXuqlGAOUO4LeYBWS+u9pSRQG0TkYfsnWj68capzaGxLUm"
    "H0BNtBMwd0EAZQ8gdgHN2/OScVsLYDZ8gfZM1/ofGsJXQFenkg/gYHQg9NyCAMofwE8/xNsTTt3blpeM"
    "21IAQ9+E8T9368OnbjQCIbswlHgAw9Gpx+weAVQggJ8efvJprGtjNh/PVkkAjWGTf+XgxkZROPRy5rPE"
    "A+iMnku2XQCVCOCnh1c+eCvOBL1L2aoIIH0x2gPIvXsg/K3hnWgV0H14KOkAQoy5GQFUJICfr1/5e7wC"
    "vspVQwBfzEQTj09lft/HCSMQbkwIYLcF8LiAf8QqYParup0fQH4qOqCVO/vMr2v0HI5mre9uertsAl0V"
    "QIUC+Pn65/EKGFut2/EBHDgSTVt77ZnluDAbRjyNjCUcQG+0E9xySgCVCuCX65//K1YB52/v9AAa5qPR"
    "z3VHjj779/r+cEWG0VKrgLcHKxLAvMOglQ/gSQFxLhlW+OzyDg+g/v1o0o7nF/SGg2Hxbi/13cjJuVcI"
    "oPudrb4x3wyW/n6CAMoVwC8fxSsg/97tHR1AoWuuxKZO556waVRTYhUw2vEKAfRNbvWNCddrbDopgAoG"
    "8HPMAtKHVndyAEML0ZSDd4denLGwiX+qt/i0t5piBDC7sSW1xfdlaGMw3IQAKhjAjz/GLKDn0PKeiphp"
    "evMBZGbD4c6llxav5nCdoLaLxS+SdTJcIWC+9OtUGFsrPZQtljAcuu7MAQFUNIC464D8veaKmB958wHk"
    "70QTnjtc/9JH78kw5LW/uejEh8LF1E8XNtnHju7z7ektvi8PwhdirvpCTGUDeLoOeGv7PMHG9jcfwGIY"
    "8nB5qsjm+9LGGbKiPxo2fym6fbL0b4rlT8YYU7GZ6RvRFtDcvpQAKhzAf7dVAeUI4OtwkPJskV+GbJ0M"
    "V2T4sujiNRtWSWdLX0e0587/PZ2wufthG218TACVDuDRo49inhHboQF0hq34by9mih0hCr8XcLPoVv7G"
    "caLV0tdSbw0D7Za/39JML4ZBt9mRHgFUPoBH17dPAWUI4FaYbE/Rqy8PPIhOdOc+LvbTwfkT0cLZXfIS"
    "WpnhsAX/bnor85y+EA61Dt5KCSCBAB7FHRm3EwOoD0dw2+4X3Y0tbHxXfn9XsSHio+G4VMmzvA331+Jd"
    "YqWEqfADnXX7pwWQSAD/uX7lf+ydzU8iyxbA1TSvmw7w0IahQW2fF1Q+lAhoQOReoAHDxwvTj7BwcKGz"
    "Ny8BE3eSMEsWQ+JqEpO3m7iazWxmZvOi+8lseZNZ3O3N/QPe+qldBQ1UfzBPELXOTqyu/qjzqzpV59Qp"
    "TfmCHiMA2yAMTp84l5mBbsF6V2nEv086JpLczkl69/+aA6c6yXEd+yQG4EEA+O8fr7TlC3p8AHRcTDqn"
    "zDKOwdfZGONH8QGfyCrXv3+GWyY2Ln/ijYtGuBArMwBgAMYBwB/XGndJPjYAZmC0z5FfrkgMRgpRTUQn"
    "z7dgDWn0LOAc8uP9YviJ/t+4130vEgPwYAD8uNa2U/6RAcAlYNLBkqx2crWOp/gQZQPBSYQ1gcogRBdg"
    "tPhueejXNZwkurlnqjI2FgZgLABMBgH3DUAE6vZ7hSV6P9wtIKbM6hPGCE2cuergOpG7kztGOb0Kevrx"
    "abO741Qnl6YXAzAeAH5cv/nLUwOArMJFTJdCwAfXhCFxAYSKkbXOgSFzgf6oaX+gkzupPjPkOoJ7+sgi"
    "OSbVR2IAHhaA/0wAAfcMQAba55seRUMcunutEYS/l6l2ummr+Z00Wo1d6J51TUVDQ7wnSZ9spSnpVruss"
    "rA8BAzAmAH78ef3vvz8pAPiwAyYAVeydiQXYjxdQZn65LtFUW3jnMEQQROh8Jmrr/q7fFVA1s++2EBIt"
    "XPUl26D25X1oGICxAfDn9Zt/kk8IAHsDxinnVEYKGBJn9SBMJXK/97QEvbfxa8Pbky1Db0bvhTE1tOxw"
    "0C87FYYPDMD4APj9+s1vf30yABi2gG1PBVRiGO8qTR0u0CTqDXkHWEVBdVXoTDQQghgqhKlyMSnDDSef"
    "aY6JePsonZPmNRt3M+gB8cYq9U0kFK9etKWHh/3AU5a0P19BrurR//HTa+hSAXhFK6Tos0Ewade9goB9"
    "vIsqXf6nCEupX0zg37Yz2KCX7hbChlCn3fu4hxY7DGWJILrs89TOlwwR7X9Sx+Mluvve4xoIz7d0dw30"
    "SLCg8zp4X3Z1sS3EEw+oWdpzW7xW8to6Ku8dV4KydxwZU/FMC8Gu4ow/zq8Kp+OOWfrWBgsdmSTcjU2g"
    "m4PXU5uytWPgC81i+FYaqnqBmQasMTviic6y0IwrBhBjZFuoYJ6EnZjKE6s44tijYIl6uHTVuebbx1R3"
    "l718/xX9yQn5LaGHASnr7CVyAFXl6nFfYYZcb7vn+TulJDAXsm+slFVlbbyF7asYfEiEkvPoF4hWcwst"
    "SOXZXr0Pufmi0A7I6XE8pXMV1rhaSerP3r2wVubHYGBh2OLdKOvP/BoKuv+XfCXLuhrzyYcM7dNlSKJb"
    "vMS6+LF7lb7CNNm+wVGvGcsdR9V+36UhbUhHVRr+5IirzvWXES5Tg1dtaGPBuqNZRLc3t9xegZ83LoCX"
    "S1NnZ95xZvO3Wwbf9is1zkq8ZOja56kF53t5YNtXGcTaLLeyTC6XlnHe2PDQ9x62BSn2t1IUPX9ywLGq"
    "Fa1stVvz5enHZee12y2XGv0bYEXDG/ao3tBz5cQupU1IqGi8W6AZ1/mYICAO1gWAyAUAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAMB/"
    "8n+xJtMt7jn4mQAAAABJRU5ErkJggg=="
)
DEFAULT_LOGO_DATA_URI = f"data:image/png;base64,{DEFAULT_LOGO_B64}"

from typing import List, Dict, Optional, Union

class ExcelRequest(BaseModel):
    titulo: str
    headers: List[str]
    rows: List[List]
    formulas: Optional[Dict[str, Union[str, List[str]]]] = None
    hojas: Optional[List[Union[Dict[str, List[List]], str]]] = None

class ExcelData(BaseModel):
    headers: List[str]
    rows: List[List[Any]]

class ExcelOptions(BaseModel):
    quality: Optional[str] = None
    theme: Optional[Dict[str, Any]] = None
    sheets: Optional[List[Dict[str, Any]]] = None
    print: Optional[Dict[str, Any]] = None

class ExcelRequestV2(BaseModel):
    titulo: Optional[str] = None
    data: ExcelData
    options: Optional[ExcelOptions] = None


class WordRequest(BaseModel):
    # Modo legado (compat)
    titulo: Optional[str] = None
    secciones: Optional[List[str]] = None
    tablas: Optional[List[List[List[str]]]] = None

    # Modo avanzado
    template_id: Optional[str] = None  
    placeholders: Optional[Dict[str, Any]] = None
    content: Optional[List[Dict[str, Any]]] = None
    options: Optional[Dict[str, Any]] = None





class PPTBrand(BaseModel):
    primary: Optional[str] = "#112B49"
    secondary: Optional[str] = "#E6EEF8"
    accent: Optional[str] = "#F5A623"
    title_font: Optional[str] = "Calibri Light"
    body_font: Optional[str] = "Calibri"
    logo_b64: Optional[str] = None
    logo_url: Optional[str] = None


# Acepta 'title' como string o lista (compat) y bullets como lista
class PowerPointSlide(BaseModel):
    title: Union[str, List[str]] = "Slide"
    bullets: Optional[List[str]] = []

# Acepta cada item de 'slides' como objeto PowerPointSlide O como string (título suelto)
class PowerPointRequest(BaseModel):
    template_id: Optional[str] = None
    title: Optional[str] = None
    subtitle: Optional[str] = None
    theme: Optional[Dict[str, Any]] = None
    options: Optional[Dict[str, Any]] = None

    titulo: Optional[str] = None                      # compat con antiguo
    bullets: Optional[List[str]] = []                 # compat
    slides: Optional[List[Union[Dict[str, Any], PowerPointSlide, str]]] = None
    apply_branding: Optional[bool] = True
    brand: Optional[PPTBrand] = None
    background: Optional[str] = None






class PDFRequest(BaseModel):
    # Legado (compat)
    titulo: Optional[str] = None
    contenido: Optional[List[str]] = None
    incluir_grafico: Optional[bool] = False
    # Avanzado
    template_id: Optional[str] = None
    title: Optional[str] = None
    meta: Optional[Dict[str, Any]] = None
    brand: Optional[Dict[str, Any]] = None  # {logo_url|logo_b64, primary}
    sections: Optional[List[Dict[str, Any]]] = None
    options: Optional[Dict[str, Any]] = None  # {page_size, footer_text, toc}


class CanvaRequest(BaseModel):
    # modo legado (compatible)
    titulo: Optional[str] = None
    elementos: Optional[List[str]] = None
    # modo avanzado
    title: Optional[str] = None
    theme: Optional[Dict[str, Any]] = None   # {bg, card, primary, text}
    kpis: Optional[List[Dict[str, Any]]] = None  # [{label, value}]
    items: Optional[List[str]] = None
    size: Optional[Dict[str, int]] = None     # {w, h}
    to_png: Optional[bool] = False


class PowerBIRequest(BaseModel):
    headers: List[str]
    rows: List[List]

class TrainModelRequest(BaseModel):
    features: List[List[float]]
    labels: List[float]

class PredictModelRequest(BaseModel):
    features: List[List[float]]

def _hex_to_rgb(hexstr: str):
    if not hexstr:
        return RGBColor(0, 0, 0)
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _style_title(title_shape, brand: PPTBrand):
    if not title_shape:
        return
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.font.name = brand.title_font or "Calibri Light"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(brand.primary or "#112B49")

def _style_body(text_frame, brand: PPTBrand):
    if not text_frame:
        return
    text_frame.word_wrap = True
    for p in text_frame.paragraphs:
        p.font.name = brand.body_font or "Calibri"
        p.font.size = Pt(20)

def _set_background(slide, color_hex: Optional[str]):
    if not color_hex:
        return
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(color_hex)

def _add_logo(slide, prs, brand: PPTBrand):
    try:
        logo_bytes = _load_logo_bytes(getattr(brand, "logo_url", None), getattr(brand, "logo_b64", None))
        if not logo_bytes:
            return
        stream = io.BytesIO(logo_bytes)
        stream.seek(0)
        picture = slide.shapes.add_picture(stream, 0, 0)
        target_height = Inches(0.9)
        scale = target_height / picture.height
        picture.height = target_height
        picture.width = int(picture.width * scale)
        margin = Inches(0.4)
        picture.left = prs.slide_width - picture.width - margin
        picture.top = Inches(0.3)
    except Exception:
        pass

def _brand_slide(slide, prs, brand: PPTBrand, company_name: Optional[str]):
    _add_logo(slide, prs, brand)
    if not company_name:
        return
    try:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(4.0), Inches(0.35))
        p = box.text_frame.paragraphs[0]
        p.text = company_name
        p.font.name = brand.body_font or "Calibri"
        p.font.size = Pt(14)
        p.font.color.rgb = _hex_to_rgb(brand.primary or "#112B49")
    except Exception:
        pass

def _add_footer(slide, prs, idx: int, total: int, brand: PPTBrand, date_text: Optional[str] = None, company_name: Optional[str] = None, show_slide_number: bool = True):
    # izquierda (fecha)
    if date_text:
        box_left = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.5), Inches(3), Inches(0.3))
        pL = box_left.text_frame.paragraphs[0]
        pL.text = date_text
        pL.font.name = brand.body_font or "Calibri"
        pL.font.size = Pt(10)
        pL.font.color.rgb = _hex_to_rgb(brand.primary or "#112B49")
    # centro (company)
    if company_name:
        box_center = slide.shapes.add_textbox(prs.slide_width / 2 - Inches(1.5), prs.slide_height - Inches(0.5), Inches(3.0), Inches(0.3))
        pC = box_center.text_frame.paragraphs[0]
        pC.text = company_name
        pC.font.name = brand.body_font or "Calibri"
        pC.font.size = Pt(10)
        pC.font.color.rgb = _hex_to_rgb(brand.primary or "#112B49")
    # derecha (X / N)
    if show_slide_number:
        box = slide.shapes.add_textbox(prs.slide_width - Inches(1.2), prs.slide_height - Inches(0.5), Inches(1.0), Inches(0.3))
        p = box.text_frame.paragraphs[0]
        p.text = f"{idx + 1} / {total}"
        p.font.name = brand.body_font or "Calibri"
        p.font.size = Pt(10)
        p.font.color.rgb = _hex_to_rgb(brand.primary or "#112B49")


def _add_slide_number(slide, prs, idx: int, brand: PPTBrand):
    box = slide.shapes.add_textbox(prs.slide_width - Inches(1.0), prs.slide_height - Inches(0.5), Inches(0.8), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = str(idx + 1)
    p.font.name = brand.body_font or "Calibri"
    p.font.size = Pt(10)
    p.font.color.rgb = _hex_to_rgb(brand.primary or "#112B49")

def _add_simple_field(paragraph, field_instr: str):
    """Inserta un campo simple (ej. PAGE, NUMPAGES, TOC...)"""
    fld = OxmlElement('w:fldSimple')
    fld.set(qn('w:instr'), field_instr)
    r = OxmlElement('w:r')
    t = OxmlElement('w:t')
    t.text = ""  # Word calculará el valor
    r.append(t)
    fld.append(r)
    paragraph._p.append(fld)

def _add_page_numbering(paragraph, pattern: str = "Página {PAGE} de {NUMPAGES}"):
    """Inserta la numeración usando campos. Acepta con o sin llaves."""
    # normaliza tokens con o sin { }
    pattern = pattern.replace("{PAGE}", "PAGE").replace("{NUMPAGES}", "NUMPAGES")
    parts = pattern.split("PAGE")
    if len(parts) == 1:
        paragraph.add_run(pattern)
        return
    # texto antes de PAGE
    paragraph.add_run(parts[0])
    _add_simple_field(paragraph, "PAGE")
    tail = "PAGE".join(parts[1:])
    parts2 = tail.split("NUMPAGES")
    paragraph.add_run(parts2[0])
    _add_simple_field(paragraph, "NUMPAGES")
    if len(parts2) > 1:
        paragraph.add_run(parts2[1])

def _clear_section_container(container):
    for tbl in list(container.tables):
        tbl._element.getparent().remove(tbl._element)
    for p in list(container.paragraphs):
        p._element.getparent().remove(p._element)

def _set_header_footer(section, header_cfg: Optional[Dict[str, str]], footer_cfg: Optional[Dict[str, str]], logo_url=None, logo_b64=None, watermark_text=None):
    # Encabezado (tres zonas simuladas con alineación)
    header = section.header
    _clear_section_container(header)
    if header_cfg:
        p = header.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
        left = header_cfg.get("left", "")
        center = header_cfg.get("center", "")
        right = header_cfg.get("right", "Página {PAGE} de {NUMPAGES}")
        # Left
        if left:
            run = p.add_run(left)
        # Center
        if center:
            pc = header.add_paragraph()
            pc.alignment = WD_ALIGN_PARAGRAPH.CENTER
            pc.add_run(center)
        # Right con numeración
        if right:
            pr = header.add_paragraph()
            pr.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            _add_page_numbering(pr, right)

    # Logo (URL o base64) al encabezado, alineado a la derecha
    logo_bytes = None
    if logo_url or logo_b64:
        logo_bytes = _load_logo_bytes(logo_url, logo_b64)
    if logo_bytes:
        p_logo = header.add_paragraph()
        p_logo.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run_logo = p_logo.add_run()
        _add_docx_image(run_logo, logo_bytes, width_in=1.6)

    if watermark_text:
        pw = header.add_paragraph()
        pw.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = pw.add_run(watermark_text)
        run.font.size = DocxPt(48)
        run.font.color.rgb = DocxRGBColor(0xB4, 0xB4, 0xB4)

    # Pie
    if footer_cfg:
        footer = section.footer
        _clear_section_container(footer)
        p = footer.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        center = footer_cfg.get("center", "")
        left = footer_cfg.get("left", "")
        right = footer_cfg.get("right", "")
        if left:
            pl = footer.add_paragraph(); pl.alignment = WD_ALIGN_PARAGRAPH.LEFT; pl.add_run(left)
        if center:
            pc = footer.add_paragraph(); pc.alignment = WD_ALIGN_PARAGRAPH.CENTER; pc.add_run(center)
        if right:
            pr = footer.add_paragraph(); pr.alignment = WD_ALIGN_PARAGRAPH.RIGHT; pr.add_run(right)

def _insert_toc(doc):
    p = doc.add_paragraph()
    _add_simple_field(p, r'TOC \o "1-3" \h \z \u')  # Word lo actualizará al abrir

def _apply_section_orientation(section, orientation: str):
    if orientation.lower() == "landscape":
        section.orientation = WD_ORIENT.LANDSCAPE
        # intercambiar ancho/alto
        new_width, new_height = section.page_height, section.page_width
        section.page_width, section.page_height = new_width, new_height
    else:
        section.orientation = WD_ORIENT.PORTRAIT

def _render_table(doc, spec: Dict[str, Any]):
    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    style = spec.get("style", "Light Grid Accent 5")
    cols = max(1, len(headers))
    tbl = doc.add_table(rows=1, cols=cols)
    tbl.style = style
    # header row
    for i, h in enumerate(headers):
        tbl.rows[0].cells[i].text = str(h)
    for row in rows:
        cells = tbl.add_row().cells
        for i, val in enumerate(row[:cols]):
            cells[i].text = str(val)
    return tbl

def _to_data_uri(img: Optional[str]) -> Optional[str]:
    """Convierte URL/base64 a data URI. Si ya es data:... lo deja igual."""
    if not img:
        return None
    if img.startswith("data:"):
        return img
    if img.startswith("http://") or img.startswith("https://"):
        try:
            with urllib.request.urlopen(img, timeout=8) as resp:
                mime = resp.headers.get_content_type() or "image/png"
                b = resp.read()
                b64 = base64.b64encode(b).decode("ascii")
                return f"data:{mime};base64,{b64}"
        except Exception:
            return img  # WeasyPrint también puede resolver URLs si hay red
    return img

def _prepare_pdf_payload(payload: Dict[str, Any]) -> Dict[str, Any]:
    """Normaliza brand/logo y secciones (ids para TOC y data URIs)."""
    pl = dict(payload)
    brand = (pl.get("brand") or {})
    # logo_url <- preferimos data URI
    logo_b64 = brand.get("logo_b64")
    logo_url = brand.get("logo_url") or DEFAULT_LOGO_URL
    if logo_b64:
        pl["logo_url"] = f"data:image/png;base64,{logo_b64}"
    else:
        pl["logo_url"] = _to_data_uri(logo_url)
    company_name = brand.get("company_name") or pl.get("company_name") or DEFAULT_COMPANY_NAME
    pl["company_name"] = company_name
    meta = dict(pl.get("meta") or {})
    meta.setdefault("company_name", company_name)
    pl["meta"] = meta

    # ids para headings + convertir imágenes a data URI
    sections = []
    hcount = 0
    for s in (pl.get("sections") or []):
        s = dict(s)
        if s.get("type") in ("h1", "h2"):
            hcount += 1
            s["_id"] = f"h{hcount}"
        if s.get("type") == "img" and s.get("src"):
            s["src"] = _to_data_uri(s["src"])
        sections.append(s)
    pl["sections"] = sections
    pl["headings"] = [s for s in sections if s.get("type") in ("h1", "h2")]
    return pl

SVG_TMPL = """<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}" viewBox="0 0 {w} {h}">
  <defs>
    <style>
      .title {{ font: 700 36px Inter, Arial, sans-serif; fill: {text}; }}
      .kpi-label {{ font: 600 14px Inter, Arial, sans-serif; fill: {primary}; letter-spacing: .5px; }}
      .kpi-value {{ font: 700 26px Inter, Arial, sans-serif; fill: {text}; }}
      .item {{ font: 500 16px Inter, Arial, sans-serif; fill: {text}; }}
    </style>
  </defs>
  <rect x="0" y="0" width="{w}" height="{h}" rx="24" fill="{bg}"/>
  <text x="48" y="80" class="title">{title}</text>

  <!-- KPIs -->
  {kpi_cards}

  <!-- Lista -->
  {bullet_list}
</svg>"""

def _safe_svg_text(s: Optional[str]) -> str:
    s = s or ""
    return s.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")

def _build_svg_panel(payload: Dict[str, Any], to_png: bool = False):
    size = payload.get("size") or {}
    w = int(size.get("w", 1200))
    h = int(size.get("h", 675))

    theme = payload.get("theme") or {}
    bg      = theme.get("bg", "#0B1220")
    card    = theme.get("card", "#111827")
    primary = theme.get("primary", "#22D3EE")
    text    = theme.get("text", "#E5E7EB")

    kpis = (payload.get("kpis") or [])[:3]
    cards = []
    col_w = (w - 96) / 3.0
    for i, k in enumerate(kpis):
        x = 48 + i * col_w
        cards.append(
            f'<g>'
            f'<rect x="{x:.0f}" y="120" width="{col_w-16:.0f}" height="120" rx="16" fill="{card}"/>'
            f'<text x="{x+20:.0f}" y="165" class="kpi-label">{_safe_svg_text(str(k.get("label","")))}</text>'
            f'<text x="{x+20:.0f}" y="200" class="kpi-value">{_safe_svg_text(str(k.get("value","")))}</text>'
            f'</g>'
        )
    kpi_cards = "\n".join(cards)

    # Lista de items (hasta 8 líneas)
    items = (payload.get("items") or [])[:8]
    y0 = 280
    lines = []
    for i, it in enumerate(items):
        y = y0 + i * 30
        lines.append(
            f'<circle cx="64" cy="{y}" r="4" fill="{primary}"/>'
            f'<text x="80" y="{y+5}" class="item">{_safe_svg_text(str(it))}</text>'
        )
    bullet_list = "\n".join(lines)

    svg = SVG_TMPL.format(
        w=w, h=h,
        bg=bg, text=text, primary=primary,
        title=_safe_svg_text(payload.get("title") or "Panel"),
        kpi_cards=kpi_cards,
        bullet_list=bullet_list
    )

    png_bytes = None
    if to_png:
        if cairosvg is None:
            raise HTTPException(
                status_code=500,
                detail='CairoSVG no está instalado. Agrega "cairosvg" a requirements.txt e instala.'
            )
        png_bytes = cairosvg.svg2png(bytestring=svg.encode("utf-8"), output_width=w, output_height=h)
    return svg, png_bytes

def _coerce_json(value: Any):
    """Convierte cadenas a objetos Python/JSON. Acepta comillas simples (literal_eval)."""
    if isinstance(value, str):
        # Primero intenta JSON puro
        try:
            return json.loads(value)
        except Exception:
            pass
        # Luego intento con literal_eval (permite comillas simples estilo Python)
        try:
            return ast.literal_eval(value)
        except Exception:
            return value
    return value

PDF_HTML_TMPL = r"""
<!doctype html>
<html>
<head>
<meta charset="utf-8">
<style>
  @page {
    size: {{ page_size }};
    margin: 18mm 16mm 20mm 16mm;
    @bottom-center {
      content: "{{ footer_text }}" " · Pág. " counter(page) " de " counter(pages);
      font-size: 10pt; color: #666;
    }
  }
  body { font-family: Inter, Arial, sans-serif; color:#111; }
  h1, h2 { color: {{ primary }}; margin: 18px 0 8px; }
  h1 { font-size: 22pt; }
  h2 { font-size: 16pt; }
  p { font-size: 11.5pt; line-height: 1.45; margin: 8px 0 10px; }
  table { width:100%; border-collapse: collapse; margin: 10px 0 14px; font-size: 11pt; }
  th, td { border:1px solid #ddd; padding: 8px; }
  th { background: {{ primary }}10; text-align: left; }
  .cover { display:flex; height:85vh; align-items:center; justify-content:center; text-align:center; }
  .logo { max-height:60px; margin-bottom: 10px; }
  .company { font-size: 12pt; color:#333; margin-top: 6px; }
  figure { margin: 10px 0; text-align:center; }
  figcaption { font-size: 10pt; color:#666; }
  /* TOC */
  .toc { margin: 12px 0 24px; }
  .toc h2 { color:#666; margin: 0 0 6px; font-size: 12pt; }
  .toc a { text-decoration:none; color:#111; display:block; padding:2px 0; }
  .toc a::after {
    content: target-counter(attr(href), page);
    float: right; color:#888;
  }
</style>
</head>
<body>

<!-- Portada -->
<div class="cover">
  <div>
    {% if logo_url %}<img class="logo" src="{{ logo_url }}">{% endif %}
    <div class="company">{{ company_name }}</div>
    <h1>{{ title }}</h1>
    <div>{{ meta.autor or "" }}{% if meta.autor and meta.fecha %} · {% endif %}{{ meta.fecha or "" }}</div>
  </div>
</div>

<!-- TOC simple (opcional) -->
{% if toc and headings|length %}
<div class="toc">
  <h2>Contenido</h2>
  {% for h in headings %}
    <a href="#{{ h._id }}">{{ h.text }}</a>
  {% endfor %}
</div>
{% endif %}

<!-- Secciones -->
{% for s in sections %}
  {% if s.type == 'h1' %}<h1 id="{{ s._id }}">{{ s.text }}</h1>{% endif %}
  {% if s.type == 'h2' %}<h2 id="{{ s._id }}">{{ s.text }}</h2>{% endif %}
  {% if s.type == 'p'  %}<p>{{ s.text }}</p>{% endif %}

  {% if s.type == 'table' %}
    <table>
      <thead><tr>{% for h in s.headers %}<th>{{ h }}</th>{% endfor %}</tr></thead>
      <tbody>
        {% for r in s.rows %}<tr>{% for c in r %}<td>{{ c }}</td>{% endfor %}</tr>{% endfor %}
      </tbody>
    </table>
  {% endif %}

  {% if s.type == 'img' %}
    <figure>
      <img src="{{ s.src }}" style="max-width:100%">
      {% if s.caption %}<figcaption>{{ s.caption }}</figcaption>{% endif %}
    </figure>
  {% endif %}
{% endfor %}

</body></html>
"""

@app.get("/healthz", include_in_schema=False)
def healthz():
    return {"ok": True}


@app.get("/")
def root():
    return {"message": "API funcionando correctamente"}

def _brand_excel_sheet(ws, max_cols: int):
    max_cols = max(1, max_cols)
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max_cols)
    cell = ws.cell(row=1, column=1)
    cell.value = DEFAULT_COMPANY_NAME
    cell.hyperlink = DEFAULT_LOGO_URL
    cell.font = Font(bold=True, size=14, color="0563C1")
    cell.alignment = Alignment(horizontal="center")
    try:
        logo_bytes = _load_logo_bytes(DEFAULT_LOGO_URL, DEFAULT_LOGO_B64)
        if logo_bytes:
            bio = io.BytesIO(logo_bytes)
            bio.seek(0)
            img = XLImage(bio)
            img.width = 120
            img.height = 120
            anchor_col = min(max_cols, 3)
            ws.add_image(img, f"{get_column_letter(anchor_col)}1")
    except Exception:
        pass

def _apply_excel_header_footer(ws):
    header = getattr(ws, "header_footer", None)
    if header is not None:
        if hasattr(header.left_header, "text"):
            header.left_header.text = DEFAULT_COMPANY_NAME
        else:
            header.left_header = DEFAULT_COMPANY_NAME
        if hasattr(header.center_header, "text"):
            header.center_header.text = ""
        else:
            header.center_header = ""
        if hasattr(header.left_footer, "text"):
            header.left_footer.text = DEFAULT_COMPANY_NAME
        else:
            header.left_footer = DEFAULT_COMPANY_NAME
        return
    # Fallback para versiones viejas de openpyxl
    odd_header = getattr(ws, "oddHeader", None)
    if odd_header is not None:
        if hasattr(odd_header.left, "text"):
            odd_header.left.text = DEFAULT_COMPANY_NAME
        else:
            odd_header.left = DEFAULT_COMPANY_NAME
        if hasattr(odd_header.center, "text"):
            odd_header.center.text = ""
        else:
            odd_header.center = ""
    even_header = getattr(ws, "evenHeader", None)
    if even_header is not None:
        if hasattr(even_header.left, "text"):
            even_header.left.text = DEFAULT_COMPANY_NAME
        else:
            even_header.left = DEFAULT_COMPANY_NAME
        if hasattr(even_header.center, "text"):
            even_header.center.text = ""
        else:
            even_header.center = ""
    odd_footer = getattr(ws, "oddFooter", None)
    if odd_footer is not None:
        if hasattr(odd_footer.left, "text"):
            odd_footer.left.text = DEFAULT_COMPANY_NAME
        else:
            odd_footer.left = DEFAULT_COMPANY_NAME
    even_footer = getattr(ws, "evenFooter", None)
    if even_footer is not None:
        if hasattr(even_footer.left, "text"):
            even_footer.left.text = DEFAULT_COMPANY_NAME
        else:
            even_footer.left = DEFAULT_COMPANY_NAME

@app.post("/generate_excel")
def generate_excel(request: Request, data: Union[ExcelRequestV2, ExcelRequest]):
    """
    Acepta:
    - v1: { "titulo": "...", "headers": [...], "rows": [...] }
    - v2: { "titulo": "...", "data": {...}, "options": {...} }
    """
    # ---- Normaliza a (titulo, headers, rows, opts) ----
    if isinstance(data, ExcelRequest):
        d = sanitize(data.dict())
        titulo = d["titulo"]
        headers = d["headers"]
        rows = d.get("rows") or []
        opts: Dict[str, Any] = {}
    else:
        # ExcelRequestV2
        d = data.dict()  # no sanitizo fórmulas ni formatos
        titulo = d.get("titulo") or "Libro"
        headers = d["data"]["headers"]
        rows = d["data"].get("rows") or []
        opts = d.get("options") or {}

    # --- Lee overrides v2 (si vinieron) ---
    theme          = (opts.get("theme") or {})
    number_formats = (theme.get("number_formats") or {})  # {"Columna": "0.00"} o formatos de Excel
    sheets_opts    = (opts.get("sheets") or [])
    s0             = sheets_opts[0] if sheets_opts else {}
    freeze_addr    = s0.get("freeze")
    widths_map     = (s0.get("widths") or {})             # {"A": 18, "B": 30, ...}
    table_opts     = (s0.get("table") or {})
    table_style    = table_opts.get("style", "Table Style Medium 9")
    totals_row     = table_opts.get("totals_row", True)

    print_opts     = (opts.get("print") or {})
    print_orient   = (print_opts.get("orientation") or "landscape").lower()
    fit_to_width   = int(print_opts.get("fit_to_width", 1))


    wb = Workbook()
    wb.remove(wb.active)

    header_fill = PatternFill("solid", fgColor="E6EEF8")
    header_font = Font(bold=True, color="112B49")

    # ====== Hoja: Detalle ======
    ws_det = wb.create_sheet("Detalle")
    _brand_excel_sheet(ws_det, max(len(headers), 2))
    _apply_excel_header_footer(ws_det)
    ws_det.append(headers)
    header_row_idx = ws_det.max_row
    data_start_idx = header_row_idx + 1
    for r in rows:
        ws_det.append(r)

    # Estilo encabezado
    for c in ws_det[header_row_idx]:
        c.fill = header_fill
        c.font = header_font

    if freeze_addr:
        ws_det.freeze_panes = freeze_addr
    else:
        ws_det.freeze_panes = f"A{data_start_idx}"

    # Ajuste de anchos auto
    for col_idx, h in enumerate(headers, start=1):
        width = max(len(str(h)), *(len(str(r[col_idx-1])) if col_idx-1 < len(r) else 0 for r in rows))
        ws_det.column_dimensions[get_column_letter(col_idx)].width = max(10, min(width + 2, 40))
    # Overrides de widths si llegaron
    for col_letter, w in widths_map.items():
        try:
            ws_det.column_dimensions[col_letter].width = float(w)
        except Exception:
            pass

    # Detectar columnas numéricas
    def is_number(v):
        try:
            float(v); return True
        except Exception:
            return False

    numeric_cols = set()
    for j in range(len(headers)):
        if any(is_number(r[j]) for r in rows if j < len(r)):
            numeric_cols.add(j)

    currency_headers = {"venta", "ventas", "costo", "costos", "margen", "importe", "monto", "total"}
    percent_headers  = {"%", "porcentaje", "ratio", "margen %"}
    from openpyxl.styles.numbers import FORMAT_CURRENCY_USD_SIMPLE, FORMAT_PERCENTAGE_00
    for j, h in enumerate(headers):
        lch = str(h).strip().lower()
        col_letter = get_column_letter(j+1)
        rng = f"{col_letter}{data_start_idx}:{col_letter}{ws_det.max_row}"
        if any(k in lch for k in currency_headers) and j in numeric_cols:
            for cell in ws_det[rng]:
                for c in cell: c.number_format = FORMAT_CURRENCY_USD_SIMPLE
        elif any(k in lch for k in percent_headers) and j in numeric_cols:
            for cell in ws_det[rng]:
                for c in cell: c.number_format = FORMAT_PERCENTAGE_00

    # Overrides de formatos por nombre de columna (v2 theme.number_formats)
    for col_name, fmt in number_formats.items():
        if col_name in headers:
            j = headers.index(col_name) + 1
            rng = f"{get_column_letter(j)}{data_start_idx}:{get_column_letter(j)}{ws_det.max_row}"
            for cell in ws_det[rng]:
                for c in cell: c.number_format = fmt

    # Tabla con o sin totales (según options)
    last_row = ws_det.max_row
    last_col = ws_det.max_column
    if totals_row:
        ws_det.append([None] * last_col)
        last_row_tot = ws_det.max_row
        ref = f"A{header_row_idx}:{get_column_letter(last_col)}{last_row_tot}"
        table = Table(displayName="TablaDetalle", ref=ref, totalsRowCount=1)
        style = TableStyleInfo(name=table_style, showRowStripes=True, showColumnStripes=False)
        table.tableStyleInfo = style
        for idx, col in enumerate(table.tableColumns, start=1):
            if idx == 1:
                col.totalsRowLabel = "Totales"
            elif (idx-1) in numeric_cols:
                col.totalsRowFunction = "sum"
        ws_det.add_table(table)
    else:
        ref = f"A{header_row_idx}:{get_column_letter(last_col)}{last_row}"
        table = Table(displayName="TablaDetalle", ref=ref)
        table.tableStyleInfo = TableStyleInfo(name=table_style, showRowStripes=True, showColumnStripes=False)
        ws_det.add_table(table)

    last_row = ws_det.max_row

    # Formato condicional sobre la última numérica (si existe)
    num_cols_sorted = sorted(list(numeric_cols))
    if num_cols_sorted:
        last_num_col = num_cols_sorted[-1] + 1
        rng = f"{get_column_letter(last_num_col)}{data_start_idx}:{get_column_letter(last_num_col)}{last_row}"
        ws_det.conditional_formatting.add(
            rng,
            ColorScaleRule(start_type="min", mid_type="percentile", mid_value=50, end_type="max")
        )

    # Validación de datos simple sobre col A (si aplica)
    colA_vals = [str(r[0]) for r in rows if len(r) > 0]
    uniques = sorted(set(colA_vals))
    if 1 <= len(uniques) <= 20 and sum(len(u) for u in uniques) < 240:
        dv = DataValidation(type="list", formula1='"' + ",".join(uniques) + '"', allow_blank=True)
        ws_det.add_data_validation(dv)
        dv.add(f"A{data_start_idx}:A{last_row}")

    # Config impresión (override con options.print si vino)
    for ws in [ws_det]:
        ws.page_setup.orientation = "portrait" if print_orient == "portrait" else "landscape"
        ws.page_setup.fitToWidth = fit_to_width
        ws.page_setup.fitToHeight = 0
        ws.print_title_rows = f"1:{header_row_idx}"

    # ====== Hoja: Resumen ======
    ws_res = wb.create_sheet("Resumen")
    _apply_excel_header_footer(ws_res)
    ws_res.append(headers)
    for r in rows:
        ws_res.append(r)
    for c in ws_res[1]:
        c.fill = header_fill
        c.font = header_font
    ws_res.freeze_panes = "A2"

    # Ajuste de anchos (auto + override)
    for col_idx, h in enumerate(headers, start=1):
        width = max(len(str(h)), *(len(str(r[col_idx-1])) if col_idx-1 < len(r) else 0 for r in rows))
        ws_res.column_dimensions[get_column_letter(col_idx)].width = max(10, min(width + 2, 40))
    for col_letter, w in widths_map.items():
        try:
            ws_res.column_dimensions[col_letter].width = float(w)
        except Exception:
            pass

    # Tabla en Resumen
    ref_res = f"A1:{get_column_letter(len(headers))}{ws_res.max_row}"
    table_res = Table(displayName="TablaResumen", ref=ref_res)
    table_res.tableStyleInfo = TableStyleInfo(name=table_style, showRowStripes=True)
    ws_res.add_table(table_res)

    # Rango con nombre
    dn = DefinedName(
        name="Datos_Resumen",
        attr_text=f"'{ws_res.title}'!$A$1:${get_column_letter(len(headers))}${ws_res.max_row}"
    )
    if hasattr(wb.defined_names, "add"):
        wb.defined_names.add(dn)
    else:
        wb.defined_names.append(dn)

    try:
        import pandas as pd
        df = pd.DataFrame(rows, columns=headers)
        for col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="ignore")
        first_col = headers[0]
        num_cols = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]
        if num_cols:
            piv = df.groupby(first_col, dropna=False)[num_cols].sum().reset_index()
            ws_piv = wb.create_sheet("Pivot")
            _apply_excel_header_footer(ws_piv)
            ws_piv.append([first_col] + num_cols)
            for row in piv.itertuples(index=False):
                ws_piv.append(list(row))
            for c in ws_piv[1]:
                c.fill = header_fill
                c.font = header_font
            ws_piv.freeze_panes = "A2"
            ref_piv = f"A1:{get_column_letter(ws_piv.max_column)}{ws_piv.max_row}"
            t_piv = Table(displayName="TablaPivot", ref=ref_piv)
            t_piv.tableStyleInfo = TableStyleInfo(name=table_style, showRowStripes=True)
            ws_piv.add_table(t_piv)
            ws_piv.page_setup.orientation = "portrait" if print_orient == "portrait" else "landscape"
            ws_piv.page_setup.fitToWidth = fit_to_width
            ws_piv.page_setup.fitToHeight = 0
            ws_piv.print_title_rows = "1:1"
    except Exception:
        pass

    # ====== Gráficos ======
    ws_chart = wb.create_sheet("Gráficos")
    _apply_excel_header_footer(ws_chart)
    ws_chart["A1"] = DEFAULT_COMPANY_NAME
    ws_chart["A1"].hyperlink = DEFAULT_LOGO_URL
    ws_chart["A1"].font = Font(bold=True, color="0563C1")
    ws_chart["A1"].alignment = Alignment(horizontal="center")
    cats = None
    if ws_res.max_row >= 2:
        cats = Reference(ws_res, min_col=1, min_row=2, max_row=ws_res.max_row)

    if len(headers) >= 2 and 1 in numeric_cols and cats is not None:
        vals = Reference(ws_res, min_col=2, min_row=1, max_row=ws_res.max_row)
        chart = BarChart()
        chart.title = "Serie principal por categoría"
        chart.add_data(vals, titles_from_data=True)
        chart.set_categories(cats)
        chart.y_axis.title = headers[1]
        chart.x_axis.title = headers[0]
        ws_chart.add_chart(chart, "A2")

    num_cols_sorted = sorted(list(numeric_cols))
    if num_cols_sorted and cats is not None:
        last_nc = num_cols_sorted[-1] + 1
        if last_nc != 2:
            vals2 = Reference(ws_res, min_col=last_nc, min_row=1, max_row=ws_res.max_row)
            chart2 = LineChart()
            chart2.title = f"{headers[last_nc-1]} (tendencia)"
            chart2.add_data(vals2, titles_from_data=True)
            chart2.set_categories(cats)
            ws_chart.add_chart(chart2, "J2")

    ws_chart.page_setup.orientation = "portrait" if print_orient == "portrait" else "landscape"
    ws_chart.page_setup.fitToWidth = fit_to_width
    ws_chart.page_setup.fitToHeight = 0

    # === Guardar ===
    safe_title = re.sub(r"[^\w\-]+", "_", titulo).strip("_") or "archivo"
    file_id = f"{safe_title}_{uuid.uuid4().hex[:8]}.xlsx"
    file_path = os.path.join(RESULT_DIR, file_id)
    wb.save(file_path)
    return {"url": _result_url(file_id, request)}

@app.post("/generate_word")
def generate_word(data: WordRequest):
    if data.content or data.placeholders or data.options or data.template_id:
        placeholders = data.placeholders or {}
        options = data.options or {}
        content = data.content or []

        doc = Document()

        # === Portada (si hay placeholders) ===
        titulo = placeholders.get("titulo") or "Documento"
        subtitulo = placeholders.get("subtitulo") or ""
        autor = placeholders.get("autor") or ""
        fecha = placeholders.get("fecha") or ""

        ptitle = doc.add_paragraph()
        ptitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = ptitle.add_run(titulo); r.bold = True; r.font.size = DocxPt(24)

        if subtitulo:
            ps = doc.add_paragraph()
            ps.alignment = WD_ALIGN_PARAGRAPH.CENTER
            rs = ps.add_run(subtitulo); rs.font.size = DocxPt(14)

        meta = []
        if autor: meta.append(autor)
        if fecha: meta.append(fecha)
        if meta:
            pm = doc.add_paragraph()
            pm.alignment = WD_ALIGN_PARAGRAPH.CENTER
            pm.add_run(" – ".join(meta)).italic = True

        doc.add_page_break()

        # === TOC opcional ===
        if options.get("toc", False):
            _insert_toc(doc)
            doc.add_page_break()

        # === Encabezado/Pie + Logo/Watermark en TODAS las secciones ===
        logo_url = placeholders.get("logo_url")
        logo_b64 = placeholders.get("logo_b64")
        wm = None
        wm_cfg = options.get("watermark")
        if isinstance(wm_cfg, dict):
            wm = wm_cfg.get("text")
        _set_header_footer(
            doc.sections[0],
            options.get("header", {"right": "Página {PAGE} de {NUMPAGES}"}),
            options.get("footer", {"center": ""}),
            logo_url=logo_url, logo_b64=logo_b64, watermark_text=wm
        )

        # Mapeo "from": "table:1", "heading:2", etc.
        sec_specs = options.get("sections", []) or []
        # contador por tipo
        counters = {"heading": 0, "paragraph": 0, "table": 0, "list": 0, "image": 0}
        for item in content:
            typ = item.get("type", "paragraph")
            for s in sec_specs:
                src = s.get("from")
                if src and ":" in src:
                    t, n = src.split(":", 1)
                    try:
                        n = int(n)
                    except Exception:
                        n = None
                    if t == typ and n == counters.get(typ, 0) + 1:
                        # nueva sección (página nueva) con orientación indicada
                        new_sec = doc.add_section(WD_SECTION_START.NEW_PAGE)
                        _apply_section_orientation(new_sec, s.get("orientation", "portrait"))
                        # heredar header/footer
                        _set_header_footer(
                            new_sec,
                            options.get("header", {"right": "Página {PAGE} de {NUMPAGES}"}),
                            options.get("footer", {"center": ""}),
                            logo_url=logo_url, logo_b64=logo_b64, watermark_text=wm
                        )
                        break

            # ahora insertamos el elemento
            if typ == "heading":
                level = int(item.get("level", 1))
                text = str(item.get("text", ""))
                para = doc.add_paragraph(text, style=f"Heading {min(max(level,1),3)}")
                counters["heading"] += 1

            elif typ == "paragraph":
                text = str(item.get("text", ""))
                para = doc.add_paragraph(text, style="Normal")
                counters["paragraph"] += 1

            elif typ == "table":
                _render_table(doc, item)
                counters["table"] += 1

            elif typ == "list":
                items = item.get("items", [])
                ordered = bool(item.get("ordered", False))
                style = "List Number" if ordered else "List Bullet"
                for it in items:
                    p = doc.add_paragraph(str(it), style=style)
                counters["list"] += 1

            elif typ == "image":
                # admite url o base64
                width_in = float(item.get("width_in", 5))
                if item.get("image_b64"):
                    try:
                        img = io.BytesIO(b64decode(item["image_b64"]))
                        doc.add_picture(img, width=DocxInches(width_in))
                    except Exception:
                        pass
                elif item.get("url") and (item["url"].startswith("http://") or item["url"].startswith("https://")):
                    try:
                        with urllib.request.urlopen(item["url"]) as resp:
                            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                            tmp.write(resp.read()); tmp.flush()
                            doc.add_picture(tmp.name, width=DocxInches(width_in))
                    except Exception:
                        pass
                counters["image"] += 1

            else:
                # fallback
                doc.add_paragraph(str(item))
        
        # === Guardar ===
        file_id = f"{uuid.uuid4()}.docx"
        file_path = os.path.join(RESULT_DIR, file_id)
        doc.save(file_path)
        return {"url": f"/resultados/{file_id}"}

    data = sanitize(data.dict())  # aquí sí sanitizamos como antes
    doc = Document()
    doc.add_heading(data["titulo"], 0)
    for sec in data["secciones"]:
        doc.add_paragraph(sec)
    if data.get("tablas"):
        for tabla in data["tablas"]:
            t = doc.add_table(rows=1, cols=len(tabla[0]))
            hdr_cells = t.rows[0].cells
            for i, h in enumerate(tabla[0]):
                hdr_cells[i].text = h
            for row in tabla[1:]:
                row_cells = t.add_row().cells
                for i, cell in enumerate(row):
                    row_cells[i].text = cell
    file_id = f"{uuid.uuid4()}.docx"
    file_path = os.path.join(RESULT_DIR, file_id)
    doc.save(file_path)
    return {"url": f"/resultados/{file_id}"}



@app.post("/generate_ppt")
def generate_ppt(request: Request, data: PowerPointRequest):
    # Modo AVANZADO si hay 'type' en los slides o si trae 'title/subtitle/theme/options/template_id'
    advanced = bool(
        (data.slides and any(isinstance(s, dict) and "type" in s for s in data.slides))
        or data.title or data.subtitle or data.theme or data.options or data.template_id
    )

    if advanced:
        # NO sanitizamos (para no romper hex, URLs, etc.)
        payload = data.dict()

        # === plantilla (mapping opcional) ===
        template_map = {
            # "corporate-v1": "templates/corporate_v1.pptx",
        }
        template_path = template_map.get(payload.get("template_id") or "", None)
        prs = Presentation(template_path) if template_path else Presentation()

        # === brand / theme ===
        brand_data = payload.get("brand") or {}
        brand = PPTBrand(**brand_data)
        if not (brand.logo_b64 or brand.logo_url):
            brand.logo_b64 = DEFAULT_LOGO_B64
            brand.logo_url = DEFAULT_LOGO_URL
        company_name = (
            brand_data.get("company_name")
            or (payload.get("company") if isinstance(payload.get("company"), str) else None)
            or payload.get("company_name")
            or DEFAULT_COMPANY_NAME
        )
        theme = payload.get("theme") or {}
        if theme.get("primary"):
            brand.primary = theme["primary"]
        if theme.get("font"):
            brand.title_font = theme["font"]
            brand.body_font  = theme["font"]

        # color de fondo global
        global_bg = payload.get("background") or brand.secondary

        # Portada si existe un slide de tipo cover (en el orden del payload)
        slides_in = payload.get("slides") or []
        for s in slides_in:
            if not isinstance(s, dict) or s.get("type") != "cover":
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title
            if global_bg: _set_background(slide, global_bg)
            # Title / Subtitle
            slide.shapes.title.text = payload.get("title") or payload.get("titulo") or "Presentación"
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = payload.get("subtitle", "") or ""
            # estilos
            if brand.primary:
                _style_title(slide.shapes.title, brand)
            _brand_slide(slide, prs, brand, company_name)

        # Resto de slides en orden
        for s in slides_in:
            if isinstance(s, dict) and s.get("type") == "cover":
                continue  # ya hecho

            # === KPIs ===
            if isinstance(s, dict) and s.get("type") == "kpis":
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
                if global_bg: _set_background(slide, global_bg)
                slide.shapes.title.text = s.get("title") or "KPIs"
                _style_title(slide.shapes.title, brand)
                # caja de KPIs
                tx = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(8.0), Inches(3.0))
                tf = tx.text_frame
                tf.clear()
                items = s.get("items", [])
                for i, item in enumerate(items):
                    p = tf.add_paragraph() if i else tf.paragraphs[0]
                    p.text = f'{item.get("label","")}: {item.get("value","")}'
                    p.level = 0
                    p.font.name = brand.body_font or "Calibri"
                    p.font.size = Pt(28)
                    p.font.color.rgb = _hex_to_rgb(brand.primary or "#112B49")
                _brand_slide(slide, prs, brand, company_name)

            # === TABLA ===
            elif isinstance(s, dict) and s.get("type") == "table":
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
                if global_bg: _set_background(slide, global_bg)
                slide.shapes.title.text = s.get("title") or "Tabla"
                _style_title(slide.shapes.title, brand)

                headers = s.get("headers", [])
                rows = s.get("rows", [])
                rows_n = len(rows) + 1
                cols_n = max(1, len(headers))
                table = slide.shapes.add_table(
                    rows_n, cols_n, Inches(0.8), Inches(1.6),
                    Inches(8.4), Inches(0.8 + rows_n * 0.35)
                ).table

                # encabezados
                for j, h in enumerate(headers):
                    cell = table.cell(0, j)
                    cell.text = str(h)
                    ph = cell.text_frame.paragraphs[0]
                    ph.font.bold = True
                    ph.font.name = brand.body_font or "Calibri"
                    ph.font.size = Pt(14)
                    ph.font.color.rgb = _hex_to_rgb(brand.primary or "#112B49")
                # filas
                for i, row in enumerate(rows, start=1):
                    for j, val in enumerate(row[:cols_n]):
                        cell = table.cell(i, j)
                        cell.text = str(val)
                        p = cell.text_frame.paragraphs[0]
                        p.font.name = brand.body_font or "Calibri"
                        p.font.size = Pt(12)
                _brand_slide(slide, prs, brand, company_name)

            # === CHART ===
            elif isinstance(s, dict) and s.get("type") == "chart":
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
                if global_bg: _set_background(slide, global_bg)
                slide.shapes.title.text = s.get("title") or "Gráfico"
                _style_title(slide.shapes.title, brand)

                data = CategoryChartData()
                data.categories = s.get("categories", [])
                for serie in s.get("series", []):
                    data.add_series(serie.get("name","Serie"), serie.get("values", []))

                slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    Inches(1), Inches(1.6), Inches(8), Inches(4),
                    data
                )
                _brand_slide(slide, prs, brand, company_name)

            # === fallback: slide texto simple (por compatibilidad)
            else:
                # soporta strings o dicts sin 'type' (antiguo)
                if isinstance(s, str):
                    s = {"title": s, "bullets": []}
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
                if global_bg: _set_background(slide, global_bg)
                slide.shapes.title.text = s.get("title", "Slide")
                _style_title(slide.shapes.title, brand)
                body = slide.placeholders[1].text_frame
                body.clear()
                for b in (s.get("bullets", []) or []):
                    p = body.add_paragraph()
                    p.text = b
                    p.level = 0
                    p.font.name = brand.body_font or "Calibri"
                    p.font.size = Pt(20)
                _style_body(body, brand)
                _brand_slide(slide, prs, brand, company_name)

        # Footer (números + fecha)
        total = len(prs.slides)
        show_nums = (payload.get("options") or {}).get("slide_numbers", True)
        footer_date = date.today().strftime("%Y-%m-%d")
        for i, sl in enumerate(prs.slides):
            _add_footer(
                sl,
                prs,
                i,
                total,
                brand,
                date_text=footer_date,
                company_name=company_name,
                show_slide_number=bool(show_nums),
            )

        file_id = f"{uuid.uuid4()}.pptx"
        file_path = os.path.join(RESULT_DIR, file_id)
        prs.save(file_path)
        return {"url": _result_url(file_id, request)}

    data = sanitize(data.dict())  # aquí sí podemos sanitizar
    apply_branding = data.get("apply_branding", True)

    # normaliza/instancia brand
    brand_data = data.get("brand") or {}
    brand = PPTBrand(**brand_data)
    if not (brand.logo_b64 or brand.logo_url):
        brand.logo_b64 = DEFAULT_LOGO_B64
        brand.logo_url = DEFAULT_LOGO_URL
    company_name = brand_data.get("company_name") or data.get("company_name") or DEFAULT_COMPANY_NAME

    prs = Presentation()

    # normalizar title si llega como lista / strings sueltos como slides
    slides_norm = []
    for s in data.get("slides", []) or []:
        if isinstance(s, str):
            slides_norm.append({"title": s, "bullets": []})
        elif isinstance(s, dict):
            t = s.get("title")
            if isinstance(t, list):
                s["title"] = " ".join(map(str, t))
            if s.get("bullets") is None:
                s["bullets"] = []
            slides_norm.append(s)
    if slides_norm:
        data["slides"] = slides_norm

    if data.get("slides"):
        for i, s in enumerate(data["slides"]):
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
            if apply_branding:
                _set_background(slide, data.get("background") or brand.secondary)
            slide.shapes.title.text = s.get("title", "Slide")
            if apply_branding:
                _style_title(slide.shapes.title, brand)

            body = slide.placeholders[1].text_frame
            body.clear()
            for bullet in s.get("bullets", []):
                p = body.add_paragraph()
                p.text = bullet
                if apply_branding:
                    p.level = 0
                    p.font.name = brand.body_font
                    p.font.size = Pt(20)

            if apply_branding:
                _brand_slide(slide, prs, brand, company_name)
                _style_body(body, brand)
        # números en modo legado
        if apply_branding:
            total = len(prs.slides)
            for i, sl in enumerate(prs.slides):
                _add_footer(
                    sl,
                    prs,
                    i,
                    total,
                    brand,
                    date_text=date.today().strftime("%Y-%m-%d"),
                    company_name=company_name,
                    show_slide_number=True,
                )

    else:
        # slide de título con bullets
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if apply_branding:
            _set_background(slide, data.get("background") or brand.secondary)
        slide.shapes.title.text = data.get("titulo") or "Presentación"
        if apply_branding:
            _style_title(slide.shapes.title, brand)

        if data.get("bullets"):
            body = slide.placeholders[1].text_frame
            body.clear()
            for b in data["bullets"]:
                p = body.add_paragraph()
                p.text = b
                if apply_branding:
                    p.level = 0
                    p.font.name = brand.body_font
                    p.font.size = Pt(20)
            if apply_branding:
                _style_body(body, brand)
        if apply_branding:
            _brand_slide(slide, prs, brand, company_name)
            _add_footer(
                slide,
                prs,
                0,
                1,
                brand,
                date_text=date.today().strftime("%Y-%m-%d"),
                company_name=company_name,
                show_slide_number=True,
            )

    file_id = f"{uuid.uuid4()}.pptx"
    file_path = os.path.join(RESULT_DIR, file_id)
    prs.save(file_path)
    return {"url": _result_url(file_id, request)}



@app.post("/generate_pdf")
def generate_pdf(request: Request, data: PDFRequest):
    # ====== MODO AVANZADO (HTML+CSS con WeasyPrint) ======
    if data.sections or data.brand or data.title or data.template_id or data.options:
        if HTML is None:
            raise HTTPException(
                status_code=500,
                detail='WeasyPrint no está instalado. Agrega "weasyprint" a requirements.txt e instala.'
            )

        # No sanitizar para no romper data URIs, hex, etc.
        payload = data.dict()

        # --- Compatibilidad: si no hay sections, convertir 'contenido' o 'content' a párrafos ---
        if not payload.get("sections"):
            legacy_blocks = payload.get("contenido") or payload.get("content")
            if isinstance(legacy_blocks, str):
                # separa por párrafos en blanco (doble salto) o saltos largos
                parts = [p.strip() for p in re.split(r"\n\s*\n|(?:\r?\n){2,}", legacy_blocks) if p.strip()]
                payload["sections"] = [{"type": "p", "text": p} for p in parts]
            elif isinstance(legacy_blocks, list):
                payload["sections"] = [{"type": "p", "text": str(p)} for p in legacy_blocks]

        # Preparar payload (logo como data URI, IDs de headings, etc.)
        pl = _prepare_pdf_payload(payload)

        primary = (pl.get("brand") or {}).get("primary", "#0F766E")
        opts = pl.get("options") or {}
        page_size = opts.get("page_size", "A4")
        footer_text = (opts.get("footer_text") or DEFAULT_COMPANY_NAME)
        title = pl.get("title") or pl.get("titulo") or "Informe"

        # Renderizar HTML con Jinja2
        html = Template(PDF_HTML_TMPL).render(
            page_size=page_size,
            footer_text=footer_text,
            primary=primary,
            logo_url=pl.get("logo_url"),
            company_name=pl.get("company_name"),
            title=title,
            meta=pl.get("meta") or {},
            sections=pl.get("sections") or [],
            headings=pl.get("headings") or [],
            toc=bool(opts.get("toc", True))  # TOC activado por defecto
        )

        # Generar PDF con WeasyPrint
        pdf_bytes = HTML(string=html, base_url=".").write_pdf()

        # Guardar y retornar URL pública configurable
        file_id = f"{uuid.uuid4()}.pdf"
        file_path = os.path.join(RESULT_DIR, file_id)
        with open(file_path, "wb") as f:
            f.write(pdf_bytes)
        return {"url": _pdf_url(file_id)}

    # ====== MODO LEGADO (FPDF) ======
    data = sanitize(data.dict())
    pdf = FPDF()
    pdf.add_page()

    # Logo corporativo en portada (best-effort)
    logo_tmp = None
    try:
        with urllib.request.urlopen(DEFAULT_LOGO_URL, timeout=8) as resp:
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            tmp.write(resp.read())
            tmp.flush()
            tmp.close()
            logo_tmp = tmp.name
    except Exception:
        logo_tmp = None

    if logo_tmp:
        try:
            pdf.image(logo_tmp, x=77, y=15, w=55)
        except Exception:
            pass
        finally:
            try:
                os.unlink(logo_tmp)
            except Exception:
                pass
        pdf.ln(45)
    else:
        pdf.ln(10)

    # Encabezado simple
    pdf.set_text_color(0, 0, 0)
    pdf.set_font("Arial", style="B", size=16)
    pdf.cell(200, 10, txt=DEFAULT_COMPANY_NAME, ln=True, align="C")
    pdf.ln(4)
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt=data.get("titulo") or "Informe", ln=True, align="C")

    # Cuerpo (línea por línea)
    for line in (data.get("contenido") or []):
        pdf.cell(200, 10, txt=str(line), ln=True, align='L')

    # Gráfico opcional de ejemplo
    if data.get("incluir_grafico"):
        plt.plot([1, 2, 3], [1, 4, 9])
        grafico_path = os.path.join(RESULT_DIR, "grafico.png")
        plt.savefig(grafico_path)
        try:
            pdf.image(grafico_path, x=10, y=80, w=100)
        except Exception:
            pass

    # Guardar y responder
    file_id = f"{uuid.uuid4()}.pdf"
    file_path = os.path.join(RESULT_DIR, file_id)
    pdf.output(file_path)
    return {"url": _result_url(file_id, request)}


@app.post("/generate_canva")
def generate_canva(request: Request, data: CanvaRequest):
    # --- modo avanzado (plantilla SVG y opcional PNG) ---
    if any([data.title, data.theme, data.kpis, data.items, data.size, data.to_png]):
        payload = data.dict()   # no sanear: preserva colores hex, etc.
        svg_str, png_bytes = _build_svg_panel(payload, to_png=bool(payload.get("to_png")))

        svg_id = f"{uuid.uuid4()}.svg"
        svg_path = os.path.join(RESULT_DIR, svg_id)
        with open(svg_path, "w", encoding="utf-8") as f:
            f.write(svg_str)

        resp = {"url": _result_url(svg_id, request)}
        if png_bytes:
            png_id = f"{uuid.uuid4()}.png"
            png_path = os.path.join(RESULT_DIR, png_id)
            with open(png_path, "wb") as f:
                f.write(png_bytes)
            resp["url_png"] = _result_url(png_id, request)
        return resp

    # --- modo legado (tu comportamiento anterior) ---
    data = sanitize(data.dict())
    file_id = f"{uuid.uuid4()}.svg"
    file_path = os.path.join(RESULT_DIR, file_id)
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(f"<svg xmlns='http://www.w3.org/2000/svg' width='800' height='400' viewBox='0 0 800 400'>")
        f.write(f"<text x='10' y='30' style='font:700 20px Arial'>{data.get('titulo','Canva')}</text>")
        for i, el in enumerate(data.get('elementos') or []):
            f.write(f"<text x='10' y='{60+i*22}' style='font:500 14px Arial'>{el}</text>")
        f.write("</svg>")
    return {"url": _result_url(file_id, request)}


@app.post("/generate_powerbi")
def generate_powerbi(request: Request, data: PowerBIRequest):
    data = sanitize(data.dict())
    df = pd.DataFrame(data["rows"], columns=data["headers"])
    file_id = f"{uuid.uuid4()}.csv"
    file_path = os.path.join(RESULT_DIR, file_id)
    df.to_csv(file_path, index=False)
    return {"url": _result_url(file_id, request)}

@app.post("/train_model")
def train_model(data: TrainModelRequest):
    accuracy = round(random.uniform(0.8, 0.99), 2)
    return {"status": "ok", "accuracy": accuracy}

@app.post("/predict_model")
def predict_model(data: PredictModelRequest):
    predictions = [random.randint(0, 1) for _ in data.features]
    return {"status": "ok", "predictions": predictions}

@app.get("/resultados/{filename}")
def get_file(filename: str, request: Request):
    file_path = os.path.join(RESULT_DIR, filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="Archivo no encontrado")
    response = FileResponse(file_path)
    if request.query_params.get("download") == "1":
        response.headers["Content-Disposition"] = f'attachment; filename="{filename}"'
    return response
    # ===========================
# 🔹 UPLOAD TO STORAGE (Render Ready)
# ===========================
import os, uuid, requests
from datetime import datetime

def upload_to_storage(file_path: str, public_path: str) -> str:
    """
    Sube el archivo al storage estático de Render o bucket público equivalente.
    Requiere variables:
      - PUBLIC_BASE_URL: dominio público de tu app
      - RENDER_WRITE_TOKEN: token de API con permiso de escritura
    """
    STORAGE_URL_BASE = os.getenv("PUBLIC_BASE_URL", "https://universal-artifact-generator.onrender.com").rstrip("/")
    WRITE_TOKEN = os.getenv("RENDER_WRITE_TOKEN")
    ENV = os.getenv("ENV", "prod")

    if not WRITE_TOKEN:
        raise RuntimeError("Falta RENDER_WRITE_TOKEN en entorno.")

    # Construye endpoint destino
    upload_url = f"{STORAGE_URL_BASE}/resultados/{os.path.basename(public_path)}"

    # Abre archivo binario y sube con PUT (Render acepta PUT/POST según config)
    with open(file_path, "rb") as f:
        resp = requests.put(upload_url, headers={"Authorization": f"Bearer {WRITE_TOKEN}"}, data=f)

    if resp.status_code not in (200, 201):
        raise RuntimeError(f"Error al subir archivo: {resp.status_code} - {resp.text[:200]}")

    print(f"[UPLOAD] Archivo subido correctamente → {upload_url}")
    return upload_url


def main(datos: dict) -> dict:
    """
    Genera un PDF, lo sube a Render y devuelve la URL pública.
    """
    pdf_path = generar_pdf(datos)
    file_uuid = str(uuid.uuid4())
    public_path = f"resultados/{file_uuid}.pdf"

    try:
        storage_url = upload_to_storage(pdf_path, public_path)
    except Exception as e:
        print(f"[ERROR] Fallo en la subida: {e}")
        storage_url = f"{os.getenv('PUBLIC_BASE_URL', 'https://universal-artifact-generator.onrender.com')}/error/{file_uuid}.pdf"

    print(f"[INFO] Archivo generado: {pdf_path} | URL pública: {storage_url} | {datetime.now()}")
    return {"url": storage_url, "id": file_uuid, "status": "ok"}


def generar_pdf(datos):
    """
    Genera un PDF temporal de ejemplo. Sustituir por tu generador real.
    """
    import tempfile
    path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pdf")
    with open(path, "w") as f:
        f.write("PDF generado de prueba\n")
    return path
# ===========================
# 🔹 ENDPOINT DE ESTADO GENERAL (Render / UptimeRobot)
# ===========================
@app.get("/healthz", include_in_schema=False)
def healthz():
    """
    Endpoint de salud para Render y UptimeRobot.
    Render lo usa para verificar que la app esté activa.
    """
    return {"status": "ok"}
from fastapi.responses import PlainTextResponse

@app.get("/status", include_in_schema=False, response_class=PlainTextResponse)
def status_plain():
    return "OK"  # 200 texto plano para monitores

@app.head("/status", include_in_schema=False)
def status_head():
    return PlainTextResponse(content="", status_code=200)
# ===========================
# 🔹 ENDPOINT PRINCIPAL (raíz)
# ===========================
@app.get("/", include_in_schema=False)
def root():
    """
    Página raíz para confirmar que la API está viva.
    Evita error 404 en la URL base.
    """
    return {"message": "API Universal Documentos funcionando correctamente ✅"}

# ===========================
# 🔹 LIMPIEZA AUTOMÁTICA DE ARCHIVOS
# ===========================
import time

@app.on_event("startup")
async def cleanup_old_results():
    """
    Elimina archivos antiguos (>1 hora) en /resultados
    para mantener el entorno Render limpio y evitar exceso de disco.
    """
    max_age = 60 * 60  # 1 hora
    now = time.time()

    if not os.path.exists(RESULT_DIR):
        return

    for f in os.listdir(RESULT_DIR):
        try:
            fp = os.path.join(RESULT_DIR, f)
            if os.path.isfile(fp) and now - os.path.getmtime(fp) > max_age:
                os.remove(fp)
        except Exception:
            pass
